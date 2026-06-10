import os
import re
import json
import unicodedata
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation


class LocalRecursiveTextSplitter:
    def __init__(self, chunk_size=1000, chunk_overlap=200):
        self.chunk_size = chunk_size
        self.chunk_overlap = min(chunk_overlap, max(0, chunk_size - 1))
        self.separators = ["\n\n", "\n", ". ", "! ", "? ", "; ", ", ", " "]

    def split_text(self, text):
        text = text.strip()
        if len(text) <= self.chunk_size:
            return [text] if text else []

        chunks = []
        start = 0
        while start < len(text):
            end = min(start + self.chunk_size, len(text))
            chunk_end = end

            if end < len(text):
                window = text[start:end]
                best_pos = -1
                for separator in self.separators:
                    pos = window.rfind(separator)
                    if pos > best_pos and pos > self.chunk_size * 0.45:
                        best_pos = pos + len(separator)
                if best_pos > 0:
                    chunk_end = start + best_pos

            chunk = text[start:chunk_end].strip()
            if chunk:
                chunks.append(chunk)

            if chunk_end >= len(text):
                break
            start = max(chunk_end - self.chunk_overlap, start + 1)

        return chunks


class DocumentProcessor:
    def __init__(self, chunk_size=1000, chunk_overlap=200):
        self.chunk_size = int(os.getenv("CHUNK_SIZE", chunk_size))
        self.chunk_overlap = int(os.getenv("CHUNK_OVERLAP", chunk_overlap))
        self.adaptive_chunking = os.getenv("ADAPTIVE_CHUNKING", "true").strip().lower() not in ["0", "false", "no", "off"]
        self.text_splitter = LocalRecursiveTextSplitter(self.chunk_size, self.chunk_overlap)

    def clean_text(self, text):
        text = re.sub(r"(\w)-\s*\n\s*(\w)", r"\1\2", text or "")
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        lines = [line.strip() for line in text.split("\n")]
        return "\n".join(lines).strip()

    def normalize_ascii(self, text):
        text = unicodedata.normalize("NFD", str(text or "").lower())
        text = "".join(ch for ch in text if unicodedata.category(ch) != "Mn")
        return re.sub(r"\s+", " ", text).strip()

    def parse_chapter_heading(self, line):
        clean = re.sub(r"\s+", " ", (line or "").strip())
        if not clean:
            return None

        patterns = [
            r"^(?:chapter|chuong)\s+([0-9]{1,2})\s*[:.\-]?\s*(.*)$",
            r"^([0-9]{1,2})\s+(Introduction|Overview of the UML Notation|Reliable, Scalable, and Maintainable Applications|Data Models and Query Languages)\b(.*)$",
        ]
        for pattern in patterns:
            match = re.match(pattern, clean, re.IGNORECASE)
            if not match:
                continue
            number = int(match.group(1))
            if not (1 <= number <= 40):
                return None
            title = ""
            if len(match.groups()) >= 2:
                title = re.sub(r"\s+", " ", " ".join(part for part in match.groups()[1:] if part)).strip(" .:-")
            title_norm = self.normalize_ascii(title)
            if title_norm and any(term in title_norm for term in [
                "described in section", "briefly described", "is described", "see chapter"
            ]):
                return None
            if title and not re.match(r"^[A-Za-z0-9]", title):
                return None
            return {"number": number, "title": title}
        return None

    def parse_section_heading(self, line):
        clean = re.sub(r"\s+", " ", (line or "").strip())
        match = re.match(r"^([0-9]{1,2}\.[0-9]+(?:\.[0-9]+)*)\s+(.{3,140})$", clean)
        if not match:
            return None
        title = match.group(2).strip(" .:-")
        if len(title) < 3:
            return None
        return {"number": match.group(1), "title": title}

    def detect_content_zone(self, text, page_number=None, heading=""):
        sample = self.normalize_ascii(f"{heading}\n{text[:1800]}")
        if any(term in sample for term in [
            "table of contents", "contents", "annotated table of contents", "muc luc"
        ]):
            return "toc"
        if any(term in sample for term in [
            "answers to exercises", "answer key", "solutions to exercises",
            "chapter 1 introduction 1 b", "multiple choice questions"
        ]):
            return "answer_key"
        if "references" in sample or "bibliography" in sample:
            return "references"
        if any(term in sample for term in ["preface", "acknowledgments", "acknowledgements"]):
            return "preface"
        if any(term in sample for term in ["appendix", "appendices"]):
            return "appendix"
        return "body"

    def make_section_path(self, chapter_number=0, chapter_title="", section_number="", section_title="", heading=""):
        parts = []
        if chapter_number:
            chapter_label = f"Chapter {chapter_number}"
            if chapter_title:
                chapter_label += f": {chapter_title}"
            parts.append(chapter_label)
        if section_number:
            section_label = section_number
            if section_title:
                section_label += f" {section_title}"
            parts.append(section_label)
        elif heading and heading not in parts:
            parts.append(heading)
        return " > ".join(part for part in parts if part)

    def detect_heading(self, text):
        lines = [line.strip() for line in self.clean_text(text).split("\n") if line.strip()]
        if not lines:
            return ""
        return self.detect_heading_line(lines[0])

    def detect_heading_line(self, line):
        line = re.sub(r"\s+", " ", (line or "").strip())[:180]
        if not line:
            return ""
        if re.match(r"^(chapter|chuong)\s+\d+(\s*[:.\-]\s*|\s+).+", line, re.IGNORECASE):
            return line
        if re.match(r"^(chapter|chuong)\s+\d+\s*$", line, re.IGNORECASE):
            return line
        if re.match(r"^\d+(\.\d+){0,4}\s+[A-Z][A-Za-z0-9 ,&:/()'\-]{3,120}$", line):
            return line
        if 8 <= len(line) <= 120 and line.isupper() and len(line.split()) <= 14:
            return line.title()
        return ""

    def detect_page_chapter(self, text):
        lines = [line.strip() for line in self.clean_text(text).split("\n") if line.strip()]
        for index, line in enumerate(lines[:14]):
            chapter = self.parse_chapter_heading(line)
            if chapter:
                if not chapter.get("title"):
                    for next_line in lines[index + 1:index + 4]:
                        if self.parse_chapter_heading(next_line):
                            continue
                        if 3 <= len(next_line) <= 140:
                            candidate = re.sub(r"\s+", " ", next_line).strip(" .:-")
                            candidate_norm = self.normalize_ascii(candidate)
                            if candidate and not re.match(r"^[A-Za-z0-9]", candidate):
                                return None
                            if any(term in candidate_norm for term in [
                                "described in section", "briefly described", "is described", "see chapter"
                            ]):
                                return None
                            chapter["title"] = candidate
                            break
                return chapter
        return None

    def split_text_into_units(self, text, page_number=None, slide_number=None, inherited_state=None):
        lines = [line.strip() for line in self.clean_text(text).split("\n")]
        units = []
        inherited_state = inherited_state or {}
        current_heading = inherited_state.get("heading", "")
        current_chapter_number = int(inherited_state.get("chapter_number") or 0)
        current_chapter_title = inherited_state.get("chapter_title", "")
        current_section_number = inherited_state.get("section_number", "")
        current_section_title = inherited_state.get("section_title", "")
        current_parts = []

        def flush():
            if not current_parts:
                return
            unit_text = self.clean_text("\n".join(current_parts))
            if not unit_text:
                return
            detected = self.detect_heading(unit_text)
            heading = current_heading or detected
            section_path = self.make_section_path(
                current_chapter_number,
                current_chapter_title,
                current_section_number,
                current_section_title,
                heading
            )
            units.append({
                "text": unit_text,
                "page_number": page_number,
                "slide_number": slide_number,
                "heading": heading,
                "section_path": section_path,
                "detected_title": detected or heading,
                "chapter_number": current_chapter_number,
                "chapter_title": current_chapter_title,
                "section_number": current_section_number,
                "section_title": current_section_title,
                "content_zone": self.detect_content_zone(unit_text, page_number, heading)
            })

        for line in lines:
            if not line:
                if current_parts:
                    current_parts.append("")
                continue
            chapter = self.parse_chapter_heading(line)
            if (
                chapter
                and not chapter.get("title")
                and current_chapter_number
                and int(chapter.get("number") or 0) != current_chapter_number
            ):
                chapter = None
            section = self.parse_section_heading(line)
            heading = self.detect_heading_line(line)
            if (chapter or section or heading) and current_parts:
                flush()
                current_parts = []
            if chapter:
                current_chapter_number = chapter["number"]
                current_chapter_title = chapter.get("title", "") or current_chapter_title
                current_section_number = ""
                current_section_title = ""
                current_heading = self.make_section_path(current_chapter_number, current_chapter_title)
            elif section:
                current_section_number = section["number"]
                current_section_title = section["title"]
                major = int(current_section_number.split(".", 1)[0])
                if major and current_chapter_number != major:
                    current_chapter_number = major
                    current_chapter_title = {
                        1: "Introduction",
                        2: "Overview"
                    }.get(major, current_chapter_title if current_chapter_number == major else "")
                current_heading = self.make_section_path(
                    current_chapter_number,
                    current_chapter_title,
                    current_section_number,
                    current_section_title
                )
            if heading:
                current_heading = heading
            current_parts.append(line)

        flush()
        return units, {
            "heading": current_heading,
            "chapter_number": current_chapter_number,
            "chapter_title": current_chapter_title,
            "section_number": current_section_number,
            "section_title": current_section_title
        }

    def extract_units_from_pdf(self, file_path):
        units = []
        current_state = {}
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page_number, page in enumerate(reader.pages, 1):
                text = self.clean_text(page.extract_text() or "")
                if not text:
                    continue
                page_chapter = self.detect_page_chapter(text)
                if page_chapter:
                    current_state = {
                        "heading": self.make_section_path(page_chapter["number"], page_chapter.get("title", "")),
                        "chapter_number": page_chapter["number"],
                        "chapter_title": page_chapter.get("title", ""),
                        "section_number": "",
                        "section_title": ""
                    }
                page_units, current_state = self.split_text_into_units(
                    text,
                    page_number=page_number,
                    inherited_state=current_state
                )
                units.extend(page_units)
        return units

    def extract_units_from_docx(self, file_path):
        doc = DocxDocument(file_path)
        units = []
        current_heading = ""
        current_parts = []

        def flush():
            if not current_parts:
                return
            text = self.clean_text("\n".join(current_parts))
            if not text:
                return
            detected = self.detect_heading(text)
            heading = current_heading or detected
            units.append({
                "text": text,
                "page_number": None,
                "slide_number": None,
                "heading": heading,
                "section_path": heading,
                "detected_title": detected or heading,
                "chapter_number": 0,
                "chapter_title": "",
                "section_number": "",
                "section_title": "",
                "content_zone": self.detect_content_zone(text, None, heading)
            })

        for para in doc.paragraphs:
            text = self.clean_text(para.text)
            if not text:
                continue
            style_name = (para.style.name if para.style else "").lower()
            looks_like_heading = style_name.startswith("heading") or bool(self.detect_heading_line(text))
            if looks_like_heading and current_parts:
                flush()
                current_parts = []
            if looks_like_heading:
                current_heading = text[:180]
            current_parts.append(text)

        flush()
        return units

    def extract_units_from_pptx(self, file_path):
        prs = Presentation(file_path)
        units = []
        for slide_number, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            text = self.clean_text("\n".join(slide_text))
            if not text:
                continue
            slide_units, _ = self.split_text_into_units(text, slide_number=slide_number)
            units.extend(slide_units)
        return units

    def merge_short_chunks(self, split_texts, min_chars=180):
        merged = []
        pending = ""
        for text in split_texts:
            clean = self.clean_text(text)
            if not clean:
                continue
            pending = self.clean_text(f"{pending}\n{clean}") if pending else clean
            if len(pending) >= min_chars:
                merged.append(pending)
                pending = ""

        if pending:
            if merged and len(pending) < min_chars:
                merged[-1] = self.clean_text(f"{merged[-1]}\n{pending}")
            else:
                merged.append(pending)
        return merged

    def build_chunk(self, chunk_text, unit, local_index, strategy="structured", strategy_score=0.0):
        return {
            "text": chunk_text,
            "page_number": unit.get("page_number"),
            "slide_number": unit.get("slide_number"),
            "heading": unit.get("heading") or "",
            "section_path": unit.get("section_path") or unit.get("heading") or "",
            "detected_title": unit.get("detected_title") or unit.get("heading") or "",
            "chapter_number": int(unit.get("chapter_number") or 0),
            "chapter_title": unit.get("chapter_title") or "",
            "section_number": unit.get("section_number") or "",
            "section_title": unit.get("section_title") or "",
            "content_zone": unit.get("content_zone") or "body",
            "local_index": local_index,
            "chunking_strategy": strategy,
            "chunking_score": round(float(strategy_score or 0), 4)
        }

    def split_units_structured(self, units, strategy="structured_heading", strategy_score=0.0):
        chunks = []
        for unit in units:
            split_texts = self.merge_short_chunks(self.text_splitter.split_text(unit["text"]))
            for local_index, chunk_text in enumerate(split_texts):
                chunk_text = self.clean_text(chunk_text)
                if len(chunk_text) < 20:
                    continue
                chunks.append(self.build_chunk(chunk_text, unit, local_index, strategy, strategy_score))
        return chunks

    def split_units_recursive_document(self, units, strategy="recursive_document", strategy_score=0.0):
        if not units:
            return []
        document_text = self.clean_text("\n\n".join(unit.get("text", "") for unit in units))
        split_texts = self.merge_short_chunks(self.text_splitter.split_text(document_text))
        chunks = []
        last_unit = units[0]
        search_cursor = 0
        for local_index, chunk_text in enumerate(split_texts):
            chunk_text = self.clean_text(chunk_text)
            if len(chunk_text) < 20:
                continue
            needle = chunk_text[:80]
            found_at = document_text.find(needle, search_cursor) if needle else -1
            if found_at >= 0:
                search_cursor = found_at + max(len(needle), 1)
                running = 0
                for unit in units:
                    running += len(unit.get("text", "")) + 2
                    if running >= found_at:
                        last_unit = unit
                        break
            chunks.append(self.build_chunk(chunk_text, last_unit, local_index, strategy, strategy_score))
        return chunks

    def split_units_page_aware(self, units, strategy="page_aware", strategy_score=0.0):
        grouped = {}
        order = []
        for unit in units:
            key = unit.get("page_number") or unit.get("slide_number") or f"unit:{len(order)}"
            if key not in grouped:
                grouped[key] = []
                order.append(key)
            grouped[key].append(unit)

        chunks = []
        local_index = 0
        for key in order:
            group_units = grouped[key]
            base_unit = group_units[0]
            text = self.clean_text("\n\n".join(unit.get("text", "") for unit in group_units))
            split_texts = self.merge_short_chunks(self.text_splitter.split_text(text))
            for chunk_text in split_texts:
                chunk_text = self.clean_text(chunk_text)
                if len(chunk_text) < 20:
                    continue
                chunks.append(self.build_chunk(chunk_text, base_unit, local_index, strategy, strategy_score))
                local_index += 1
        return chunks

    def chunk_lengths(self, chunks):
        return [len(chunk.get("text", "")) for chunk in chunks if chunk.get("text")]

    def score_chunk_strategy(self, chunks, units):
        if not chunks:
            return {"score": 0.0, "size": 0.0, "integrity": 0.0, "metadata": 0.0, "density": 0.0, "reasons": ["no_chunks"]}

        lengths = self.chunk_lengths(chunks)
        target_min = max(180, int(self.chunk_size * 0.45))
        target_max = int(self.chunk_size * 1.18)
        size_ok = sum(1 for length in lengths if target_min <= length <= target_max) / max(len(lengths), 1)
        too_short = sum(1 for length in lengths if length < 120) / max(len(lengths), 1)
        too_long = sum(1 for length in lengths if length > self.chunk_size * 1.35) / max(len(lengths), 1)

        heading_units = sum(1 for unit in units if unit.get("heading") or unit.get("section_path") or unit.get("chapter_number"))
        metadata_chunks = sum(1 for chunk in chunks if chunk.get("section_path") or chunk.get("chapter_number"))
        metadata_score = metadata_chunks / max(len(chunks), 1)
        if heading_units == 0:
            metadata_score = max(metadata_score, 0.65)

        broken_sentence = 0
        for chunk in chunks:
            text = chunk.get("text", "").strip()
            if not text:
                continue
            starts_bad = bool(re.match(r"^[,;:)\]]|^[a-z][a-z]+\s", text))
            ends_bad = bool(re.search(r"\b(and|or|the|a|an|of|to|for|with|in|on|by)$", text[-80:].strip(), re.IGNORECASE))
            if starts_bad or ends_bad:
                broken_sentence += 1
        integrity_score = 1 - (broken_sentence / max(len(chunks), 1))

        source_chars = sum(len(unit.get("text", "")) for unit in units)
        chunk_chars = sum(lengths)
        if source_chars <= 0:
            density_score = 0.0
        else:
            ratio = chunk_chars / source_chars
            density_score = max(0.0, 1 - abs(1.0 - min(ratio, 1.35)) / 0.35)

        score = (
            size_ok * 0.34 +
            integrity_score * 0.26 +
            metadata_score * 0.24 +
            density_score * 0.16 -
            too_short * 0.12 -
            too_long * 0.10
        )
        return {
            "score": round(max(0.0, min(score, 1.0)), 4),
            "size": round(size_ok, 4),
            "integrity": round(integrity_score, 4),
            "metadata": round(metadata_score, 4),
            "density": round(density_score, 4),
            "reasons": []
        }

    def describe_strategy_result(self, strategy, data, selected=False):
        reasons = []
        if data.get("metadata", 0) >= 0.75:
            reasons.append("keeps chapter/page metadata well")
        elif data.get("metadata", 0) < 0.35:
            reasons.append("weak chapter/page metadata")

        if data.get("integrity", 0) >= 0.85:
            reasons.append("rarely cuts sentences awkwardly")
        elif data.get("integrity", 0) < 0.65:
            reasons.append("may cut sentence boundaries")

        if data.get("size", 0) >= 0.65:
            reasons.append("balanced chunk sizes")
        elif data.get("size", 0) < 0.35:
            reasons.append("many chunks are too short or too long")

        if data.get("density", 0) >= 0.85:
            reasons.append("preserves most extracted text")
        elif data.get("density", 0) < 0.60:
            reasons.append("lower text coverage")

        if not reasons:
            reasons.append("acceptable fallback strategy")
        if selected:
            reasons.insert(0, "selected by highest score")

        return {
            "strategy": strategy,
            "score": round(float(data.get("score", 0)), 4),
            "size": round(float(data.get("size", 0)), 4),
            "integrity": round(float(data.get("integrity", 0)), 4),
            "metadata": round(float(data.get("metadata", 0)), 4),
            "density": round(float(data.get("density", 0)), 4),
            "reason": "; ".join(reasons[:4])
        }

    def build_chunking_report(self, scores, selected_strategy):
        strategies = [
            self.describe_strategy_result(name, data, selected=(name == selected_strategy))
            for name, data in sorted(scores.items(), key=lambda item: item[1].get("score", 0), reverse=True)
        ]
        selected = next((item for item in strategies if item["strategy"] == selected_strategy), strategies[0] if strategies else {})
        return {
            "enabled": bool(self.adaptive_chunking),
            "selected_strategy": selected_strategy,
            "selected_score": selected.get("score", 0),
            "summary": selected.get("reason", "selected strategy for this document"),
            "strategies": strategies
        }

    def with_strategy_report(self, chunks, strategy, report):
        score = float(report.get("selected_score") or 0)
        report_json = json.dumps(report, ensure_ascii=False, separators=(",", ":"))[:6000]
        for index, chunk in enumerate(chunks):
            chunk["chunking_strategy"] = strategy
            chunk["chunking_score"] = round(score, 4)
            chunk["chunking_report"] = report_json
            chunk["chunking_reason"] = str(report.get("summary") or "")[:500]
            chunk["local_index"] = index
        return chunks

    def split_units(self, units):
        if not self.adaptive_chunking:
            chunks = self.split_units_structured(units, "structured_heading", 1.0)
            report = {
                "enabled": False,
                "selected_strategy": "structured_heading",
                "selected_score": 1.0,
                "summary": "adaptive chunking is disabled; using structured heading splitter",
                "strategies": [
                    {
                        "strategy": "structured_heading",
                        "score": 1.0,
                        "size": 1.0,
                        "integrity": 1.0,
                        "metadata": 1.0,
                        "density": 1.0,
                        "reason": "adaptive chunking is disabled; default strategy used"
                    }
                ]
            }
            return self.with_strategy_report(chunks, "structured_heading", report)

        candidates = {
            "structured_heading": self.split_units_structured(units, "structured_heading", 0),
            "recursive_document": self.split_units_recursive_document(units, "recursive_document", 0),
            "page_aware": self.split_units_page_aware(units, "page_aware", 0)
        }
        scores = {name: self.score_chunk_strategy(chunks, units) for name, chunks in candidates.items()}

        best_name = max(scores, key=lambda name: scores[name]["score"])
        report = self.build_chunking_report(scores, best_name)
        best_score = report["selected_score"]
        chunks = self.with_strategy_report(candidates[best_name], best_name, report)

        summary = ", ".join(
            f"{name}={data['score']:.2f}" for name, data in sorted(scores.items())
        )
        print(f"  Adaptive chunking selected {best_name} ({best_score:.2f}); candidates: {summary}", flush=True)
        return chunks

    def process_file(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            units = self.extract_units_from_pdf(file_path)
        elif ext == ".docx":
            units = self.extract_units_from_docx(file_path)
        elif ext in [".pptx", ".ppt"]:
            units = self.extract_units_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file extension: {ext}")

        extracted_chars = sum(len(unit["text"]) for unit in units)
        if extracted_chars < 50:
            print(f"  Warning: Very little text extracted ({extracted_chars} chars)", flush=True)
            return []

        chunks = self.split_units(units)
        selected_strategy = chunks[0].get("chunking_strategy", "structured_heading") if chunks else "none"
        selected_score = chunks[0].get("chunking_score", 0) if chunks else 0
        print(
            f"  Extracted {extracted_chars} chars -> {len(units)} sections -> {len(chunks)} chunks "
            f"(strategy={selected_strategy}, score={selected_score}, size={self.chunk_size}, overlap={self.chunk_overlap})",
            flush=True
        )
        return chunks
